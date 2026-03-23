from django.shortcuts import render, redirect
from django.contrib.auth import login, authenticate, logout
from django.contrib.auth.forms import AuthenticationForm
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from .forms import RegisterForm,QuizForm,CourseMaterialForm
from .models import User, Semester, Batch, Course

# -------------------------
# Home Page
# -------------------------
def index(request):
    return render(request, "index.html")

# -------------------------
# Register View
# -------------------------
from django.contrib import messages
from django.contrib.auth import login
from django.shortcuts import render, redirect
from .forms import RegisterForm

def register(request):
    if request.method == "POST":
        form = RegisterForm(request.POST)
        if form.is_valid():
            user = form.save(commit=False)
            role = form.cleaned_data.get("role")

            if role == "Teacher":
                user.is_approved = False
                user.batch = None
            else:
                user.is_approved = True

            user.save()

            # ✅ Don't log in automatically — just show message and redirect to login page
            if role == "Teacher":
                messages.info(request, "✅ Your account is pending admin approval. You cannot log in yet.")
            else:
                messages.success(request, "🎉 Registration successful! Please log in to continue.")

            return redirect("login")

        else:
            messages.error(request, "Please correct the errors below.")
    else:
        form = RegisterForm()

    return render(request, "register.html", {"form": form})

# -------------------------
# Login View
# -------------------------
from django.contrib.auth import authenticate, login
from django.contrib.auth.forms import AuthenticationForm
from django.contrib import messages
from django.shortcuts import render, redirect

def login_view(request):
    if request.user.is_authenticated:
        # Already logged in users get redirected directly
        if request.user.role == "Student":
            return redirect("student_dashboard")
        elif request.user.role == "Teacher":
            return redirect("teacher_dashboard")
        else:
            return redirect("admin_dashboard")

    if request.method == "POST":
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            username = form.cleaned_data["username"]
            password = form.cleaned_data["password"]
            user = authenticate(request, username=username, password=password)

            if user is not None:
                # Check teacher approval status
                if user.role == "Teacher" and not user.is_approved:
                    messages.warning(request, "Your account is pending admin approval.")
                    return redirect("login")

                # Successful login
                login(request, user)
                messages.success(request, f"Welcome back, {user.username}!")

                # Role-based redirects
                role_redirects = {
                    "Student": "student_dashboard",
                    "Teacher": "teacher_dashboard",
                    "Admin": "admin_dashboard"
                }
                return redirect(role_redirects.get(user.role, "home"))

            else:
                messages.error(request, "Invalid username or password.")
        else:
            messages.error(request, "Please correct the errors.")
    else:
        form = AuthenticationForm()

    return render(request, "login.html", {"form": form})

# -------------------------
# Logout View
# -------------------------
from django.contrib.auth import logout
from django.contrib import messages
from django.shortcuts import redirect

def logout_view(request):
    if request.user.is_authenticated:
        username = request.user.username  # store before logout
        logout(request)
        messages.success(request, f"Goodbye, {username}! You’ve been logged out successfully.")
    else:
        messages.info(request, "You are not logged in.")

    return redirect("index")


# -------------------------
# Dashboards
# -------------------------


from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from django.contrib.auth.decorators import login_required, user_passes_test
from .models import Course, Semester, Batch, User
from django.db.models import Q

import re
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from django.contrib.auth.decorators import login_required, user_passes_test
from .models import Course, Semester, Batch, User
from django.db.models import Q

# Only admins can access
def is_admin(user):
    return user.is_authenticated and user.role == "Admin"

@login_required
@user_passes_test(is_admin)
def admin_dashboard(request):
    if request.method == "POST":
        active_tab = request.POST.get("active_tab", "courses")
        action = request.POST.get("action")

        # ---------- Courses ----------
        if action == "add_course" or "edit_course_id" in request.POST:
            course_code = request.POST.get("course_code", "").strip().upper()
            course_name = request.POST.get("course_name", "").strip()
            semester_id = request.POST.get("semester")
            teacher_id = request.POST.get("assigned_teacher")
            errors = []

            # Required fields
            if not course_code:
                errors.append("Course code is required.")
            if not course_name:
                errors.append("Course name is required.")
            if not semester_id:
                errors.append("Semester is required.")

            # Course code format
            if course_code and not re.match(r"^[A-Z]{3,5}\d{3}$", course_code):
                errors.append("Course code must be in format like SJMCA101.")

            # Validate semester
            semester = None
            if semester_id:
                try:
                    semester = Semester.objects.get(id=semester_id)
                except Semester.DoesNotExist:
                    errors.append("Selected semester does not exist.")

            # Validate teacher
            assigned_teacher = None
            if teacher_id:
                assigned_teacher = User.objects.filter(id=teacher_id, role="Teacher", is_approved=True).first()
                if not assigned_teacher:
                    errors.append("Selected teacher does not exist or is not approved.")

            # Unique course code
            course_qs = Course.objects.filter(course_code=course_code)
            if "edit_course_id" in request.POST:
                course_id = request.POST.get("edit_course_id")
                course_qs = course_qs.exclude(id=course_id)
            if course_qs.exists():
                errors.append("Course code already exists.")

            # Unique course name per semester
            if semester:
                name_qs = Course.objects.filter(name=course_name, semester=semester)
                if "edit_course_id" in request.POST:
                    name_qs = name_qs.exclude(id=request.POST.get("edit_course_id"))
                if name_qs.exists():
                    errors.append("Course name already exists in this semester.")

            # If errors, show messages
            if errors:
                for e in errors:
                    messages.error(request, e)
                return redirect(f"{request.path}?active_tab=courses")

            # CREATE / UPDATE
            if action == "add_course":
                Course.objects.create(
                    course_code=course_code,
                    name=course_name,
                    semester=semester,
                    assigned_teacher=assigned_teacher
                )
                messages.success(request, "Course added successfully!")
            elif "edit_course_id" in request.POST:
                course = get_object_or_404(Course, id=request.POST.get("edit_course_id"))
                course.course_code = course_code
                course.name = course_name
                course.semester = semester
                course.assigned_teacher = assigned_teacher
                course.save()
                messages.success(request, "Course updated successfully!")

            return redirect(f"{request.path}?active_tab=courses")

        # Delete Course
        elif "delete_course_id" in request.POST:
            course = get_object_or_404(Course, id=request.POST.get("delete_course_id"))
            # Optional: prevent deletion if students are enrolled
            if course.student_set.exists():
                messages.error(request, "Cannot delete course. Students are enrolled.")
            else:
                course.delete()
                messages.success(request, "Course deleted successfully!")
            return redirect(f"{request.path}?active_tab=courses")

        # ---------- Batches ----------
        elif action == "add_batch":
            batch_name = request.POST.get("batch_name", "").strip()
            semester_id = request.POST.get("batch_semester")
            errors = []

            if not batch_name:
                errors.append("Batch name is required.")
            if not semester_id:
                errors.append("Semester is required.")

            semester = None
            if semester_id:
                try:
                    semester = Semester.objects.get(id=semester_id)
                except Semester.DoesNotExist:
                    errors.append("Selected semester does not exist.")

            # Unique batch per semester
            if semester and Batch.objects.filter(name=batch_name, semester=semester).exists():
                errors.append("Batch name already exists in this semester.")

            if errors:
                for e in errors:
                    messages.error(request, e)
                return redirect(f"{request.path}?active_tab=batches")

            Batch.objects.create(name=batch_name, semester=semester)
            messages.success(request, "Batch added successfully!")
            return redirect(f"{request.path}?active_tab=batches")

        elif "edit_batch_id" in request.POST:
            batch = get_object_or_404(Batch, id=request.POST.get("edit_batch_id"))
            batch_name = request.POST.get("batch_name", "").strip()
            semester_id = request.POST.get("batch_semester")
            errors = []

            if not batch_name:
                errors.append("Batch name is required.")
            if not semester_id:
                errors.append("Semester is required.")

            semester = None
            if semester_id:
                try:
                    semester = Semester.objects.get(id=semester_id)
                except Semester.DoesNotExist:
                    errors.append("Selected semester does not exist.")

            if semester and Batch.objects.filter(name=batch_name, semester=semester).exclude(id=batch.id).exists():
                errors.append("Batch name already exists in this semester.")

            if errors:
                for e in errors:
                    messages.error(request, e)
                return redirect(f"{request.path}?active_tab=batches")

            batch.name = batch_name
            batch.semester = semester
            batch.save()
            messages.success(request, "Batch updated successfully!")
            return redirect(f"{request.path}?active_tab=batches")

        elif "delete_batch_id" in request.POST:
            batch = get_object_or_404(Batch, id=request.POST.get("delete_batch_id"))
            # Optional: prevent deletion if students are enrolled
            if batch.user_set.filter(role="Student").exists():
                messages.error(request, "Cannot delete batch. Students are assigned to it.")
            else:
                batch.delete()
                messages.success(request, "Batch deleted successfully!")
            return redirect(f"{request.path}?active_tab=batches")

        # ---------- Semesters ----------
        elif action == "add_semester":
            semester_name = request.POST.get("semester_name", "").strip()
            if not semester_name:
                messages.error(request, "Semester name is required.")
            elif Semester.objects.filter(name=semester_name).exists():
                messages.error(request, "Semester name already exists.")
            else:
                Semester.objects.create(name=semester_name)
                messages.success(request, "Semester added successfully!")
            return redirect(f"{request.path}?active_tab=semesters")

        elif "edit_semester_id" in request.POST:
            sem = get_object_or_404(Semester, id=request.POST.get("edit_semester_id"))
            semester_name = request.POST.get("semester_name", "").strip()
            if not semester_name:
                messages.error(request, "Semester name is required.")
            elif Semester.objects.filter(name=semester_name).exclude(id=sem.id).exists():
                messages.error(request, "Semester name already exists.")
            else:
                sem.name = semester_name
                sem.save()
                messages.success(request, "Semester updated successfully!")
            return redirect(f"{request.path}?active_tab=semesters")

        elif "delete_semester_id" in request.POST:
            sem = get_object_or_404(Semester, id=request.POST.get("delete_semester_id"))
            # Optional: prevent deletion if courses or batches exist
            if Course.objects.filter(semester=sem).exists() or Batch.objects.filter(semester=sem).exists():
                messages.error(request, "Cannot delete semester. Courses or batches exist for it.")
            else:
                sem.delete()
                messages.success(request, "Semester deleted successfully!")
            return redirect(f"{request.path}?active_tab=semesters")

        # ---------- Teacher Approvals ----------
        elif "approve" in request.POST or "delete" in request.POST:
            approve_ids = request.POST.getlist("approve")
            delete_ids = request.POST.getlist("delete")

            if approve_ids:
                User.objects.filter(id__in=approve_ids, role="Teacher").update(is_approved=True)
                messages.success(request, "Selected teachers approved!")

            if delete_ids:
                # Optional: prevent deleting teachers assigned to courses
                for tid in delete_ids:
                    teacher = User.objects.filter(id=tid, role="Teacher").first()
                    if teacher and teacher.course_set.exists():
                        messages.error(request, f"Cannot delete {teacher.username}. Assigned to courses.")
                    else:
                        teacher.delete()
                        messages.success(request, "Selected teachers deleted!")
            return redirect(f"{request.path}?active_tab=teachers")

        # ---------- Delete user ----------
        elif "delete_user_id" in request.POST:
            user_to_delete = get_object_or_404(User, id=request.POST.get("delete_user_id"))
            # Prevent deletion if student assigned to batch or teacher assigned to courses
            if user_to_delete.role == "Student" and user_to_delete.batch:
                messages.error(request, "Cannot delete student. Assigned to a batch.")
            elif user_to_delete.role == "Teacher" and user_to_delete.course_set.exists():
                messages.error(request, "Cannot delete teacher. Assigned to courses.")
            else:
                user_to_delete.delete()
                messages.success(request, "User deleted successfully!")
            return redirect(f"{request.path}?active_tab=users")

    # ---------- GET Handling ----------
    courses = Course.objects.select_related("semester", "assigned_teacher").all()
    semesters = Semester.objects.all()
    batches = Batch.objects.select_related("semester").all()

    teachers_pending = User.objects.filter(role="Teacher", is_approved=False)
    teachers_approved = User.objects.filter(role="Teacher", is_approved=True)

    # Users Tab Filtering
    filter_role = request.GET.get("filter_role", "All")
    active_tab = request.GET.get("active_tab", "courses")

    if filter_role == "Students":
        all_users = User.objects.select_related("batch").filter(role="Student")
    elif filter_role == "Teachers":
        all_users = User.objects.select_related("batch").filter(role="Teacher", is_approved=True)
    else:
        all_users = User.objects.select_related("batch").filter(
            Q(role="Student") | Q(role="Teacher", is_approved=True)
        )

    context = {
        "courses": courses,
        "semesters": semesters,
        "batches": batches,
        "teachers_pending": teachers_pending,
        "teachers_approved": teachers_approved,
        "all_users": all_users,
        "filter_role": filter_role,
        "course_count": courses.count(),
        "semester_count": semesters.count(),
        "batch_count": batches.count(),
        "teacher_count": teachers_approved.count(),
        "active_tab": active_tab,
    }

    return render(request, "admin_dashboard.html", context)

# -------------------------
# Teacher Approval (Admin only)
# -------------------------
@login_required
def approve_teachers(request):
    if request.user.role != "Admin":
        return redirect("login")

    teachers = User.objects.filter(role="Teacher", is_approved=False)

    if request.method == "POST":
        approved_ids = request.POST.getlist("approve")
        User.objects.filter(id__in=approved_ids).update(is_approved=True)
        messages.success(request, "Selected teachers have been approved.")
        return redirect("approve_teachers")

    return render(request, "approve_teachers.html", {"teachers": teachers})


from .models import User, Semester, Batch, Course

# -------------------------
# Manage Teachers (Admin)
# -------------------------
@login_required
def manage_teachers(request):
    if request.user.role != "Admin":
        return redirect("login")

    pending_teachers = User.objects.filter(role="Teacher", is_approved=False)
    approved_teachers = User.objects.filter(role="Teacher", is_approved=True)

    if request.method == "POST":
        approve_ids = request.POST.getlist("approve")
        User.objects.filter(id__in=approve_ids).update(is_approved=True)

        delete_ids = request.POST.getlist("delete")
        User.objects.filter(id__in=delete_ids).delete()

        messages.success(request, "Changes applied successfully!")
        return redirect("manage_teachers")

    return render(request, "manage_teachers.html", {"pending_teachers": pending_teachers, "approved_teachers": approved_teachers})


# -------------------------
# Admin: Semesters
# -------------------------
@login_required
def manage_semesters(request):
    if request.user.role != "Admin":
        return redirect("login")

    semesters = Semester.objects.all()
    if request.method == "POST":
        name = request.POST.get("name").strip()
        if name:
            if not Semester.objects.filter(name__iexact=name).exists():
                Semester.objects.create(name=name)
                messages.success(request, "Semester created successfully!")
            else:
                messages.error(request, "Semester already exists!")
        else:
            messages.error(request, "Semester name cannot be empty.")
        return redirect("manage_semesters")

    return render(request, "manage_semesters.html", {"semesters": semesters})

# -------------------------
# Admin: Batches
# -------------------------
@login_required
def manage_batches(request):
    if request.user.role != "Admin":
        return redirect("login")

    batches = Batch.objects.all()
    semesters = Semester.objects.all()

    if request.method == "POST":
        name = request.POST.get("name").strip()
        semester_id = request.POST.get("semester")
        if name and semester_id:
            semester = Semester.objects.get(id=semester_id)
            Batch.objects.create(name=name, semester=semester)
            messages.success(request, "Batch created successfully!")
        else:
            messages.error(request, "All fields are required.")
        return redirect("manage_batches")

    return render(request, "manage_batches.html", {"batches": batches, "semesters": semesters})

# -------------------------
# Admin: Courses
# -------------------------
from django.shortcuts import render, redirect
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from .models import Course, Semester, User,Quiz,Module
@login_required
def manage_courses(request):
    if request.user.role != "Admin":
        return redirect("login")

    courses = Course.objects.all()
    teachers = User.objects.filter(role="Teacher", is_approved=True)
    semesters = Semester.objects.all()

    if request.method == "POST":
        # Adding a new course
        course_code = request.POST.get("course_code")
        course_name = request.POST.get("course_name")
        semester_id = request.POST.get("semester")
        teacher_id = request.POST.get("teacher")

        semester = Semester.objects.get(id=semester_id)
        teacher = User.objects.get(id=teacher_id) if teacher_id else None

        # Check unique constraint
        if Course.objects.filter(course_code=course_code, semester=semester).exists():
            messages.error(request, "This course code already exists for the selected semester.")
        else:
            Course.objects.create(
                course_code=course_code,
                name=course_name,
                semester=semester,
                assigned_teacher=teacher
            )
            messages.success(request, "Course added successfully!")
        return redirect("manage_courses")

    context = {
        "courses": courses,
        "teachers": teachers,
        "semesters": semesters
    }
    return render(request, "manage_courses.html", context)


from django.contrib.auth import get_user_model
from django.contrib import messages
from django.shortcuts import render, redirect

User = get_user_model()

def simple_password_reset(request):
    if request.method == "POST":
        username = request.POST.get("username")
        new_password = request.POST.get("new_password")
        confirm_password = request.POST.get("confirm_password")

        if new_password != confirm_password:
            messages.error(request, "Passwords do not match!")
            return redirect("simple_password_reset")

        try:
            user = User.objects.get(username=username)
            user.set_password(new_password)
            user.save()
            messages.success(request, "Password reset successful! Please login.")
            return redirect("login")
        except User.DoesNotExist:
            messages.error(request, "Username does not exist.")
            return redirect("simple_password_reset")

    return render(request, "simple_password_reset.html")

from django.shortcuts import render
from django.contrib.auth.decorators import login_required
from django.utils import timezone
from .models import Course, Quiz


from django.shortcuts import render
from django.contrib.auth.decorators import login_required
from django.utils import timezone
from .models import Course, Quiz

@login_required
def teacher_dashboard(request):
    now = timezone.localtime()

    # 1️⃣ Fetch all courses assigned to this teacher
    assigned_courses = Course.objects.filter(
        assigned_teacher=request.user
    ).select_related('semester')

    # Annotate each course with students, modules, and quizzes info
    for course in assigned_courses:
        # Students in the course and semester
        students_in_semester = course.students.filter(
            role="Student",
            batch__semester=course.semester
        )
        course.students_in_semester = students_in_semester
        course.student_count = students_in_semester.count()

        # Modules count
        course.num_modules = course.modules.count()

        # Quizzes count across all modules of this course
        course.num_quizzes = sum(module.quizzes.count() for module in course.modules.all())

    assigned_course_ids = assigned_courses.values_list('id', flat=True)

    # 2️⃣ Fetch Quizzes
    live_quizzes = Quiz.objects.filter(
        course_id__in=assigned_course_ids,
        status="Published",
        start_time__lte=now,
        end_time__gte=now
    ).select_related('course', 'module')

    upcoming_quizzes = Quiz.objects.filter(
        course_id__in=assigned_course_ids,
        status="Published",
        start_time__gt=now
    ).select_related('course', 'module')

    past_quizzes = Quiz.objects.filter(
        course_id__in=assigned_course_ids,
        status="Published",
        end_time__lt=now
    ).select_related('course', 'module')

    # 3️⃣ Render context
    context = {
        'assigned_courses': assigned_courses,
        'live_quizzes': live_quizzes,
        'upcoming_quizzes': upcoming_quizzes,
        'past_quizzes': past_quizzes,
    }

    return render(request, 'teacher_dashboard.html', context)



@login_required
def quiz_detail(request, quiz_id):
    quiz = get_object_or_404(Quiz, id=quiz_id)

    # fetch all questions of the quiz
    questions = quiz.questions.all()

    return render(request, "quiz_detail.html", {
        "quiz": quiz,
        "questions": questions
    })
 



from django.shortcuts import render, get_object_or_404, redirect
from .models import Course, Module
from .forms import ModuleForm
from django.contrib.auth.decorators import login_required

@login_required
def create_module(request, course_id):
    course = get_object_or_404(Course, id=course_id)
    
    if request.method == 'POST':
        form = ModuleForm(request.POST)
        if form.is_valid():
            module = form.save(commit=False)
            module.course = course
            module.save()
            return redirect('course_detail', course_id=course.id)
    else:
        form = ModuleForm()
    
    return render(request, 'create_module.html', {'form': form, 'course': course})


@login_required
def course_detail(request, course_id):
    course = get_object_or_404(Course, id=course_id)
    modules = course.modules.all()

    show_modal = False
    modal_error = None

    if request.method == "POST" and "name" in request.POST:
        module_name = request.POST.get("name", "").strip()

        if not module_name:
            show_modal = True
            modal_error = "Module name cannot be empty."
        elif Module.objects.filter(course=course, name__iexact=module_name).exists():
            show_modal = True
            modal_error = f"A module named '{module_name}' already exists."
        else:
            Module.objects.create(course=course, name=module_name)
            return redirect("course_detail", course_id=course.id)

    context = {
        "course": course,
        "modules": modules,
        "show_modal": show_modal,
        "modal_error": modal_error,
    }
    return render(request, "course_detail.html", context)


@login_required
def edit_module(request, course_id, module_id):
    module = get_object_or_404(Module, id=module_id, course__id=course_id)
    if request.method == 'POST':
        new_name = request.POST.get('name', '').strip()
        if not new_name:
            messages.error(request, "Module name cannot be empty.")
        elif Module.objects.filter(course=module.course, name__iexact=new_name).exclude(id=module.id).exists():
            messages.error(request, f"A module named '{new_name}' already exists.")
        else:
            module.name = new_name
            module.save()
            messages.success(request, f"Module '{module.name}' updated successfully.")
    return redirect('course_detail', course_id=course_id)


@login_required
def delete_module(request, course_id, module_id):
    module = get_object_or_404(Module, id=module_id, course__id=course_id)
    if request.method == 'POST':
        module.delete()
        messages.success(request, f"Module '{module.name}' deleted successfully.")
    return redirect('course_detail', course_id=course_id)



from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Module, CourseMaterial
from .forms import CourseMaterialForm
from django.contrib.auth.decorators import login_required

@login_required
def upload_material(request, course_id, module_id):
    course = get_object_or_404(Course, id=course_id)
    module = get_object_or_404(Module, id=module_id)

    if request.method == 'POST':
        form = CourseMaterialForm(request.POST, request.FILES)
        if form.is_valid():
            material = form.save(commit=False)
            material.course = course
            material.uploaded_by = request.user
            material.save()
            return redirect('module_detail', course_id=course.id, module_id=module.id)
    else:
        form = CourseMaterialForm(initial={'module': module})

    return render(request, 'upload_material.html', {'form': form, 'course': course, 'module': module})


from django.contrib import messages
from .models import CourseMaterial

@login_required
def delete_material(request, material_id):
    material = get_object_or_404(CourseMaterial, id=material_id)
    
    # Optional: Only allow the user who uploaded to delete
    if material.uploaded_by != request.user:
        messages.error(request, "You don't have permission to delete this material.")
        return redirect(request.META.get('HTTP_REFERER', '/'))
    
    material.file.delete(save=False)  # Remove file from storage
    material.delete()  # Remove database record
    messages.success(request, "Material deleted successfully.")
    return redirect(request.META.get('HTTP_REFERER', '/'))


from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.urls import reverse
from .models import Course, Module, CourseMaterial, Quiz, Question, Batch
from .utils import generate_mcqs_with_gemini
from datetime import datetime
import os
import tempfile
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# -------------------------
# Helper Function: Extract text from uploaded file
# -------------------------
def extract_text_from_file(uploaded_file):
    """Extract text from PDF, DOCX, or PPTX uploaded file."""
    try:
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            for chunk in uploaded_file.chunks():
                tmp.write(chunk)
            tmp_path = tmp.name

        ext = os.path.splitext(uploaded_file.name)[1].lower()
        text = ""

        if ext == ".pdf":
            reader = PdfReader(tmp_path)
            text = "\n".join([page.extract_text() or "" for page in reader.pages])

        elif ext == ".docx":
            doc = Document(tmp_path)
            text = "\n".join([p.text for p in doc.paragraphs])

        elif ext == ".pptx":
            prs = Presentation(tmp_path)
            text = "\n".join([
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            ])

        os.remove(tmp_path)
        return text.strip()
    except Exception as e:
        print("⚠️ File text extraction failed:", e)
        return ""

# -------------------------
# MODULE DETAIL VIEW
# -------------------------
@login_required
def module_detail(request, course_id, module_id):
    course = get_object_or_404(Course, id=course_id)
    module = get_object_or_404(Module, id=module_id, course=course)

    # Fetch materials and quizzes
    materials = module.materials.all().order_by('-uploaded_at')
    draft_quizzes = Quiz.objects.filter(module=module, status="Draft").order_by('-created_at')
    published_quizzes = Quiz.objects.filter(module=module, status="Published").order_by('-created_at')

    # AI-generated questions grouped by topic
    generated_questions = Question.objects.filter(module=module, is_ai_generated=True).order_by('-id')
    topics = {}
    for q in generated_questions:
        topics.setdefault(q.topic_tag or "No Topic", []).append(q)

    active_tab = request.GET.get('tab', 'materials')
    batches = Batch.objects.filter(semester=course.semester)

    # -------------------------
    # Handle POST requests
    # -------------------------
    if request.method == "POST":

        # 1️⃣ Upload Material
        if 'file' in request.FILES and 'upload_material' in request.POST:
            topic = request.POST.get('topic', '')
            file = request.FILES['file']
            CourseMaterial.objects.create(course=course, module=module, topic=topic, file=file, uploaded_by=request.user)
            messages.success(request, "Material uploaded successfully!")
            return redirect(f"{reverse('module_detail', kwargs={'course_id': course.id, 'module_id': module.id})}?tab=materials")

        # 2️⃣ Generate AI Questions
        elif 'generate_ai' in request.POST:
            topic = request.POST.get("topic") or "General"
            content = request.POST.get("content", "").strip()
            uploaded_file = request.FILES.get("upload_file")  # fixed field name
            num_questions = int(request.POST.get("num_questions") or 5)

            # Extract text from uploaded file if no content
            if uploaded_file and not content:
                content = extract_text_from_file(uploaded_file)

            if not content:
                messages.warning(request, "Please enter study content or upload a file.")
                return redirect(f"{reverse('module_detail', kwargs={'course_id': course.id, 'module_id': module.id})}?tab=aiquiz")

            # Generate questions with Gemini
            ai_questions = generate_mcqs_with_gemini(content, num_questions)
            if not ai_questions:
                messages.error(request, "AI could not generate questions. Try again!")
                return redirect(f"{reverse('module_detail', kwargs={'course_id': course.id, 'module_id': module.id})}?tab=aiquiz")

            # Save AI-generated questions
            for q in ai_questions:
                text = q.get("question_text", "").strip()
                if not text:
                    continue
                Question.objects.create(
                    course=course,
                    module=module,
                    topic_tag=topic,
                    question_text=text,
                    option_a=q.get("option_a", ""),
                    option_b=q.get("option_b", ""),
                    option_c=q.get("option_c", ""),
                    option_d=q.get("option_d", ""),
                    correct_ans=q.get("correct_ans", "A"),
                    created_by=request.user,
                    is_ai_generated=True
                )

            messages.success(request, f"{len(ai_questions)} AI-generated questions saved under topic '{topic}'!")
            return redirect(f"{reverse('module_detail', kwargs={'course_id': course.id, 'module_id': module.id})}?tab=questions")

        # 3️⃣ Create Draft Quiz from Selected Questions
        elif 'create_draft_quiz' in request.POST:
            selected_ids = request.POST.getlist("question_ids")
            if not selected_ids:
                messages.warning(request, "No questions selected to create a draft quiz.")
                return redirect(f"{reverse('module_detail', kwargs={'course_id': course.id, 'module_id': module.id})}?tab=questions")

            quiz_title = request.POST.get("quiz_title") or f"Draft Quiz ({datetime.now().strftime('%d-%b-%Y %H:%M')})"

            # Create draft quiz
            draft_quiz = Quiz.objects.create(
                course=course,
                module=module,
                teacher=request.user,
                title=quiz_title,
                status="Draft"
            )

            # Assign selected questions
            selected_questions = Question.objects.filter(id__in=selected_ids)
            for q in selected_questions:
                q.quiz = draft_quiz
                q.save()

            messages.success(request, f"Draft Quiz '{draft_quiz.title}' created with {len(selected_questions)} questions!")
            return redirect(f"{reverse('module_detail', kwargs={'course_id': course.id, 'module_id': module.id})}?tab=quizzes")

    # -------------------------
    # Render GET requests
    # -------------------------
    context = {
        'course': course,
        'module': module,
        'materials': materials,
        'draft_quizzes': draft_quizzes,
        'published_quizzes': published_quizzes,
        'topics': topics,
        'active_tab': active_tab,
        'batches': batches
    }
    return render(request, 'module_detail.html', context)




# -------------------------
# Create Quiz from selected AI questions
# -------------------------
@login_required
def create_quiz_from_questions(request):
    if request.method == "POST":
        module_id = request.POST.get("module_id")
        module = get_object_or_404(Module, id=module_id)
        course = module.course
        title = request.POST.get("title") or f"Quiz - {module.name}"
        question_ids = request.POST.getlist("question_ids")
        questions = Question.objects.filter(id__in=question_ids)

        if not questions:
            messages.error(request, "No questions selected!")
            return redirect(f"{reverse('module_detail', kwargs={'course_id': course.id, 'module_id': module.id})}?tab=questions")

        quiz = Quiz.objects.create(
            course=course,
            module=module,
            teacher=request.user,
            title=title,
            status="Draft"
        )
        for q in questions:
            q.quiz = quiz
            q.save()

        messages.success(request, f"Quiz '{quiz.title}' created from AI questions!")
        return redirect(f"{reverse('module_detail', kwargs={'course_id': course.id, 'module_id': module.id})}?tab=quizzes")
    return redirect('dashboard')

@login_required
def edit_question(request, question_id):
    question = get_object_or_404(Question, id=question_id)
    if request.method == "POST":
        question.question_text = request.POST.get("question_text")
        question.option_a = request.POST.get("option_a")
        question.option_b = request.POST.get("option_b")
        question.option_c = request.POST.get("option_c")
        question.option_d = request.POST.get("option_d")
        question.correct_ans = request.POST.get("correct_ans")
        question.save()
        messages.success(request, "Question updated successfully.")
    return redirect(request.META.get('HTTP_REFERER', '/'))

@login_required
def delete_question(request, question_id):
    question = get_object_or_404(Question, id=question_id)
    module = question.module
    question.delete()
    messages.success(request, "Question deleted successfully.")
    return redirect(f"/course/{module.course.id}/module/{module.id}/?tab=questions")

@login_required
def add_to_quiz(request, question_id):
    question = get_object_or_404(Question, id=question_id)
    # Logic to choose a quiz or show a form to select existing quiz
    # For simplicity, you can redirect to quizzes list
    messages.success(request, "You can now add this question to an existing quiz.")
    return redirect(request.META.get('HTTP_REFERER', '/'))



# -------------------------
# Quiz Management
# -------------------------
@login_required
def view_quiz(request, quiz_id):
    quiz = get_object_or_404(Quiz, id=quiz_id)
    questions = quiz.questions.all()
    return render(request, "view_quiz.html", {"quiz": quiz, "questions": questions})


@login_required
def edit_quiz(request, quiz_id):
    quiz = get_object_or_404(Quiz, id=quiz_id)
    if request.method == "POST":
        quiz.title = request.POST.get("title")
        quiz.save()
        messages.success(request, "Quiz updated successfully.")
    return redirect('module_detail', course_id=quiz.course.id, module_id=quiz.module.id)


from django.shortcuts import render, redirect
from django.contrib.auth.decorators import login_required
from .models import Quiz, Course, CourseMaterial, QuizAnswer,QuizAttempt

from django.utils import timezone

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.utils import timezone
from .models import Quiz, Course, CourseMaterial, QuizAttempt, Question

# ---------------------------
# Student Dashboard
# ---------------------------
@login_required
def student_dashboard(request):
    user = request.user
    if user.role != "Student":
        return redirect("dashboard")  # Redirect non-students appropriately

    now = timezone.localtime(timezone.now())  # timezone-aware IST now
    student_batch = user.batch
    student_semester = student_batch.semester if student_batch else None

    # ----------------------------------------
    # Quizzes
    # ----------------------------------------
    if student_batch:
        # Active quizzes: Published, assigned to batch, within start-end time, not attempted
        active_quizzes = Quiz.objects.filter(
            status="Published",
            batches=student_batch,
            start_time__lte=now,
            end_time__gte=now
        ).exclude(quizattempt__student=user).order_by("start_time")

        # Upcoming quizzes: Published, assigned to batch, start time in future, not attempted
        upcoming_quizzes = Quiz.objects.filter(
            status="Published",
            batches=student_batch,
            start_time__gt=now
        ).exclude(quizattempt__student=user).order_by("start_time")

        # Missed quizzes: Published, assigned to batch, ended already, not attempted
        missed_quizzes = Quiz.objects.filter(
            status="Published",
            batches=student_batch,
            end_time__lt=now
        ).exclude(quizattempt__student=user).order_by("-end_time")
    else:
        active_quizzes = upcoming_quizzes = missed_quizzes = []

    # Past attempts
    past_attempts = QuizAttempt.objects.filter(
        student=user
    ).select_related('quiz', 'quiz__course', 'quiz__module').order_by("-submitted_at")

    # ----------------------------------------
    # Courses & Materials
    # ----------------------------------------
    if student_semester:
        courses = Course.objects.filter(semester=student_semester)
        materials = CourseMaterial.objects.filter(
            course__in=courses
        ).order_by('-uploaded_at')
    else:
        courses = materials = []

    # ----------------------------------------
    # Optional: Leaderboard (can be implemented later)
    # ----------------------------------------
    leaderboard = None

    context = {
        "student_batch": student_batch,
        "student_semester": student_semester,
        "active_quizzes": active_quizzes,
        "upcoming_quizzes": upcoming_quizzes,
        "missed_quizzes": missed_quizzes,
        "past_attempts": past_attempts,
        "courses": courses,
        "materials": materials,
        "leaderboard": leaderboard,
        "now": now,  # pass current time for template checks
    }

    return render(request, "student_dashboard.html", context)

# ---------------------------
# Attempt Quiz
#
# ---------------------------
# Quiz Result
# ---------------------------
@login_required
def quiz_result(request, attempt_id):
    attempt = get_object_or_404(QuizAttempt, id=attempt_id, student=request.user)
    quiz = attempt.quiz
    questions = quiz.questions.all()

    return render(request, "quiz_result.html", {
        "attempt": attempt,
        "quiz": quiz,
        "questions": questions,
    })

from django.shortcuts import render, redirect
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from .models import Quiz, QuizAttempt, User

@login_required
def quiz_attempts(request, quiz_id):
    try:
        quiz = Quiz.objects.get(id=quiz_id)
    except Quiz.DoesNotExist:
        messages.error(request, "Quiz not found.")
        return redirect('teacher_dashboard')

    batches = quiz.batches.all()
    students = User.objects.filter(role="Student", batch__in=batches).distinct()

    student_attempts = []

    for student in students:
        attempt = QuizAttempt.objects.filter(quiz=quiz, student=student).first()
        if attempt:
            total_marks = quiz.total_marks or 1
            percentage = float(attempt.score) / float(total_marks) * 100  # ✅ fixed
            if percentage >= 85:
                performance = "Excellent"
            elif percentage >= 70:
                performance = "Good"
            elif percentage >= 50:
                performance = "Average"
            else:
                performance = "Poor"

            student_attempts.append({
                'student': student,
                'attempted': True,
                'score': attempt.score,
                'total_marks': total_marks,
                'performance': performance,
                'attempted_on': attempt.submitted_at
            })
        else:
            student_attempts.append({
                'student': student,
                'attempted': False
            })

    context = {
        'quiz': quiz,
        'student_attempts': student_attempts
    }

    return render(request, 'quiz_attempts.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from django.utils import timezone
from .models import Quiz, Question, QuizAttempt



# ------------------------------
# Attempt Quiz (Student)
# ------------------------------
from datetime import timedelta

@login_required
def attempt_quiz(request, quiz_id):
    student = request.user

    if student.role != 'Student':
        messages.error(request, "Only students can attempt quizzes.")
        return redirect('home')

    quiz = get_object_or_404(Quiz, id=quiz_id, status='Published')

    # Get or create attempt
    attempt, created = QuizAttempt.objects.get_or_create(
        quiz=quiz,
        student=student,
        defaults={'started_at': timezone.now()}
    )

    # Calculate dynamic end_time per attempt using quiz.duration_minutes
    end_time = attempt.started_at + timedelta(minutes=quiz.duration or 0)

    # Handle form submission
    if request.method == "POST":
        for question in quiz.questions.all():
            selected_option = request.POST.get(f"q{question.id}")
            if selected_option:
                is_correct = selected_option == question.correct_ans
                QuizAnswer.objects.update_or_create(
                    attempt=attempt,
                    question=question,
                    defaults={'selected_option': selected_option, 'is_correct': is_correct}
                )

        # Calculate total score
        total_score = 0
        for ans in attempt.answers.all():
            if ans.is_correct:
                total_score += float(quiz.marks_per_question or ans.question.marks)

        attempt.score = total_score
        attempt.submitted_at = timezone.now()
        attempt.save()

        messages.success(request, "Quiz submitted successfully!")
        return redirect('quiz_result', attempt_id=attempt.id)

    # GET request: show quiz
    questions = quiz.questions.all()
    student_answers = {ans.question.id: ans.selected_option for ans in attempt.answers.all()}

    context = {
        'quiz': quiz,
        'questions': questions,
        'student_answers': student_answers,
        'attempt': attempt,
        'end_time': end_time,  # Pass per-attempt end time to template
    }

    return render(request, 'attempt_quiz.html', context)

# ------------------------------
# Submit Quiz (Student)
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.utils import timezone
from .models import Quiz, QuizAttempt, QuizAnswer

@login_required
def submit_quiz(request, quiz_id):
    student = request.user
    quiz = get_object_or_404(Quiz, id=quiz_id, status='Published')
    attempt, created = QuizAttempt.objects.get_or_create(
        quiz=quiz,
        student=student,
        defaults={'started_at': timezone.now()}
    )

    if request.method == "POST":
        total_score = 0
        for question in quiz.questions.all():
            selected_option = request.POST.get(f"q{question.id}")
            if selected_option:
                is_correct = selected_option == question.correct_ans
                if is_correct:
                    total_score += question.marks
                QuizAnswer.objects.update_or_create(
                    attempt=attempt,
                    question=question,
                    defaults={'selected_option': selected_option, 'is_correct': is_correct}
                )
        attempt.score = total_score
        attempt.submitted_at = timezone.now()
        attempt.save()

        # Redirect to results page
        return redirect('quiz_result', attempt_id=attempt.id)

    return redirect('attempt_quiz', quiz_id=quiz.id)

from django.shortcuts import render, get_object_or_404
from django.contrib.auth.decorators import login_required
from .models import Quiz, QuizAttempt

from django.contrib.auth.decorators import login_required
from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import Quiz

@login_required
def delete_quiz(request, quiz_id):
    quiz = get_object_or_404(Quiz, id=quiz_id)
    
    # Optional: Only allow teacher/admin to delete
    if request.user.is_staff or request.user == quiz.created_by.user:
        quiz.delete()
        messages.success(request, "Quiz deleted successfully.")
    else:
        messages.error(request, "You do not have permission to delete this quiz.")

    return redirect("teacher_dashboard")

from django.shortcuts import render, get_object_or_404
from django.contrib.auth.decorators import login_required
from .models import QuizAttempt, QuizAnswer

@login_required
def quiz_result(request, attempt_id):
    attempt = get_object_or_404(QuizAttempt, id=attempt_id, student=request.user)
    quiz = attempt.quiz

    question_answers = []
    correct = wrong = not_answered = 0

    for ans in attempt.answers.all():
        if ans.selected_option:
            if ans.is_correct:
                correct += 1
            else:
                wrong += 1
        else:
            not_answered += 1

        marks_obtained = ans.question.marks if ans.is_correct else 0
        question_answers.append({
            'question': ans.question,
            'student_answer': ans.selected_option,
            'is_correct': ans.is_correct,
            'marks_obtained': marks_obtained
        })

    total_questions = quiz.questions.count()
    total_score = float(attempt.score)
    max_score = total_questions * float(getattr(quiz, 'marks_per_question', 1))
    percentage = round((total_score / max_score) * 100, 2) if max_score > 0 else 0

    # Grade system (customizable)
    if percentage >= 90:
        grade = 'A+'
    elif percentage >= 80:
        grade = 'A'
    elif percentage >= 70:
        grade = 'B'
    elif percentage >= 60:
        grade = 'C'
    elif percentage >= 50:
        grade = 'D'
    else:
        grade = 'F'

    context = {
        'quiz': quiz,
        'attempt': attempt,
        'question_answers': question_answers,
        'correct': correct,
        'wrong': wrong,
        'not_answered': not_answered,
        'percentage': percentage,
        'grade': grade
    }

    return render(request, 'quiz_result.html', context)


from django.db.models import Avg, Count, F
from django.contrib.auth.decorators import login_required
from django.shortcuts import render
from evaluation.models import Quiz, QuizAttempt, Question, Course, User

@login_required
def analytics_dashboard(request):
    if request.user.role != "Teacher":
        return render(request, "unauthorized.html")

    # Get quizzes created by this teacher
    teacher_quizzes = Quiz.objects.filter(teacher=request.user)

    # Total students across teacher’s courses
    total_students = User.objects.filter(
        role="Student",
        courses_enrolled__in=teacher_quizzes.values_list("course", flat=True)
    ).distinct().count()

    # Total quizzes conducted
    total_quizzes = teacher_quizzes.count()

    # Average score (across all attempts)
    avg_score = QuizAttempt.objects.filter(
        quiz__in=teacher_quizzes
    ).aggregate(avg=Avg("score"))["avg"] or 0

    # Find topic with lowest average score (weakest topic)
    weakest_topic = (
        Question.objects.filter(quiz__in=teacher_quizzes)
        .values("topic_tag")
        .annotate(avg_score=Avg("quiz__quizattempt__score"))
        .order_by("avg_score")
        .first()
    )
    weakest_topic_name = weakest_topic["topic_tag"] if weakest_topic else "N/A"

    # Data for performance chart (average per student)
    student_performance = (
        QuizAttempt.objects.filter(quiz__in=teacher_quizzes)
        .values(name=F("student__username"))
        .annotate(avg_score=Avg("score"))
        .order_by("-avg_score")[:6]
    )

    # Question difficulty distribution
    difficulty_counts = (
        Question.objects.filter(quiz__in=teacher_quizzes)
        .values("difficulty")
        .annotate(count=Count("id"))
    )

    # Topic-wise average performance
    topic_stats = (
        Question.objects.filter(quiz__in=teacher_quizzes)
        .values("topic_tag")
        .annotate(
            avg_score=Avg("quiz__quizattempt__score"),
            question_count=Count("id")
        )
    )

    context = {
        "total_students": total_students,
        "total_quizzes": total_quizzes,
        "avg_score": round(avg_score, 2),
        "weakest_topic": weakest_topic_name,
        "student_performance": list(student_performance),
        "difficulty_data": list(difficulty_counts),
        "topic_stats": topic_stats,
    }

    return render(request, "analytics.html", context)


from django.db.models import Avg

@property
def average_score(self):
    avg = self.quizattempt_set.aggregate(Avg('score'))['score__avg']
    return round(avg or 0, 1)




from django.utils import timezone
from django.contrib import messages
from django.shortcuts import get_object_or_404, redirect
from django.contrib.auth.decorators import login_required
from datetime import datetime
from .models import Quiz

def validate_quiz_data(request, quiz):
    """Validate quiz form data. Returns (valid, cleaned_data)"""
    start_time = request.POST.get("start_time")
    end_time = request.POST.get("end_time")
    duration = request.POST.get("duration")
    total_marks = request.POST.get("total_marks")
    marks_per_question = request.POST.get("marks_per_question")
    selected_batches = request.POST.getlist("batches")

    errors = []

    # Required fields
    if not start_time or not end_time:
        errors.append("Start and End time are required.")
    if not duration or int(duration) <= 0:
        errors.append("Duration must be a positive number.")
    if not total_marks or float(total_marks) <= 0:
        errors.append("Total marks must be greater than 0.")
    if not selected_batches:
        errors.append("Select at least one batch.")
    if marks_per_question and float(marks_per_question) > float(total_marks):
        errors.append("Marks per question cannot exceed total marks.")

    # Time parsing
    try:
        start = timezone.make_aware(datetime.strptime(start_time, "%Y-%m-%dT%H:%M"))
        end = timezone.make_aware(datetime.strptime(end_time, "%Y-%m-%dT%H:%M"))
    except Exception:
        errors.append("Invalid date format.")
        start = end = None

    # Logical checks
    if start and end:
        if start >= end:
            errors.append("End time must be after Start time.")
        if start < timezone.now():
            errors.append("Start time cannot be in the past.")

    # Quiz must have questions
    if quiz.questions.count() == 0:
        errors.append("Cannot publish a quiz without questions.")

    if errors:
        for error in errors:
            messages.error(request, error)
        return False, None

    # Cleaned data
    cleaned_data = {
        "start": start,
        "end": end,
        "duration": int(duration),
        "total_marks": float(total_marks),
        "marks_per_question": float(marks_per_question) if marks_per_question else 1.0,
        "selected_batches": selected_batches,
    }
    return True, cleaned_data

# ------------------------------
# Publish Quiz (Teacher/Admin)
# ------------------------------
from django.utils import timezone
from django.contrib import messages
from django.shortcuts import get_object_or_404, redirect
from .models import Quiz
from datetime import datetime

@login_required
def publish_quiz(request, quiz_id):
    quiz = get_object_or_404(Quiz, id=quiz_id)

    if request.method == "POST":
        start_time = request.POST.get("start_time")
        end_time = request.POST.get("end_time")
        duration = request.POST.get("duration")
        total_marks = request.POST.get("total_marks")
        marks_per_question = request.POST.get("marks_per_question")
        selected_batches = request.POST.getlist("batches")

        # Validations
        if not start_time or not end_time:
            messages.error(request, "Start and End time are required.")
            return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)
        if not duration or int(duration) <= 0:
            messages.error(request, "Duration must be positive.")
            return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)
        if not total_marks or float(total_marks) <= 0:
            messages.error(request, "Total marks must be greater than 0.")
            return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)
        if not selected_batches:
            messages.error(request, "Select at least one batch.")
            return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)
        if marks_per_question and float(marks_per_question) > float(total_marks):
            messages.error(request, "Marks per question cannot exceed total marks.")
            return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)

        # Time parsing
        try:
            start = timezone.make_aware(datetime.strptime(start_time, "%Y-%m-%dT%H:%M"))
            end = timezone.make_aware(datetime.strptime(end_time, "%Y-%m-%dT%H:%M"))
        except Exception:
            messages.error(request, "Invalid date format.")
            return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)
        if start >= end:
            messages.error(request, "End time must be after start time.")
            return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)
        if start < timezone.now():
            messages.error(request, "Start time cannot be in the past.")
            return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)
        if quiz.questions.count() == 0:
            messages.error(request, "Cannot publish a quiz without questions.")
            return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)

        # --- Save quiz ---
        quiz.start_time = start
        quiz.end_time = end
        quiz.duration = duration
        quiz.total_marks = total_marks
        quiz.marks_per_question = float(marks_per_question) if marks_per_question else 1.0
        quiz.status = "Published"
        quiz.batches.set(selected_batches)
        quiz.save()

        # --- Update each question's marks ---
        for question in quiz.questions.all():
            question.marks = quiz.marks_per_question
            question.save()

        messages.success(request, "Quiz published successfully!")
        return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)

    return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)


from django.shortcuts import get_object_or_404, redirect, render
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.utils import timezone
from datetime import datetime
from .models import Quiz, Batch

@login_required
def republish_quiz(request, quiz_id):
    quiz = get_object_or_404(Quiz, id=quiz_id)

    # ✅ Only batches from the same semester as the course
    batches = Batch.objects.filter(semester=quiz.course.semester)

    if request.method == "POST":
        start_time = request.POST.get("start_time")
        end_time = request.POST.get("end_time")
        duration = request.POST.get("duration")
        total_marks = request.POST.get("total_marks")
        marks_per_question = request.POST.get("marks_per_question")
        selected_batches = request.POST.getlist("batches")

        error_messages = []
        now = timezone.now()

        # -------------------- Validation --------------------
        if not start_time or not end_time:
            error_messages.append("Start and End time are required.")
        else:
            try:
                start = timezone.make_aware(datetime.strptime(start_time, "%Y-%m-%dT%H:%M"))
                end = timezone.make_aware(datetime.strptime(end_time, "%Y-%m-%dT%H:%M"))
                if start >= end:
                    error_messages.append("End time must be after start time.")
                if start < now:
                    error_messages.append("Start time cannot be in the past.")
                if end <= now:
                    error_messages.append("End time must be in the future.")
            except Exception:
                error_messages.append("Invalid date format.")

        if not duration or not duration.isdigit() or int(duration) <= 0:
            error_messages.append("Duration must be a positive integer.")

        if not total_marks or float(total_marks) <= 0:
            error_messages.append("Total marks must be greater than 0.")

        if marks_per_question:
            try:
                marks_val = float(marks_per_question)
                if marks_val > float(total_marks):
                    error_messages.append("Marks per question cannot exceed total marks.")
            except ValueError:
                error_messages.append("Marks per question must be a valid number.")
        else:
            marks_val = 1.0

        if not selected_batches:
            error_messages.append("Select at least one batch.")

        if quiz.questions.count() == 0:
            error_messages.append("Cannot republish a quiz without questions.")

        # -------------------- Return errors if any --------------------
        if error_messages:
            for err in error_messages:
                messages.error(request, err)
            return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)

        # -------------------- Save changes --------------------
        quiz.start_time = start
        quiz.end_time = end
        quiz.duration = int(duration)
        quiz.total_marks = float(total_marks)
        quiz.marks_per_question = marks_val
        quiz.status = "Published"
        quiz.batches.set(selected_batches)
        quiz.save()

        # Update question marks
        for question in quiz.questions.all():
            question.marks = marks_val
            question.save()

        messages.success(request, "Quiz republished successfully!")
        return redirect("module_detail", course_id=quiz.course.id, module_id=quiz.module.id)

    # For GET requests — prefill republish form
    return render(request, "republish_quiz.html", {
        "quiz": quiz,
        "batches": batches
    })

from decimal import Decimal
from django.shortcuts import render
from django.contrib.auth.decorators import login_required
from django.db.models import Avg
from .models import Course, Quiz, QuizAttempt


@login_required
def teacher_analytics(request):
    teacher = request.user

    # 1️⃣ All courses assigned to this teacher
    courses = Course.objects.filter(assigned_teacher=teacher)

    course_data = []
    total_quizzes = 0
    total_attempts_count = 0
    total_score_sum = 0
    success_total = 0

    # 2️⃣ For each course → get its published quizzes
    for course in courses:
        quizzes = Quiz.objects.filter(course=course, status='Published')
        quiz_data = []

        for quiz in quizzes:
            attempts = QuizAttempt.objects.filter(quiz=quiz)
            total_attempts = attempts.count()
            avg_score = attempts.aggregate(avg_score=Avg('score'))['avg_score'] or 0
            success_count = attempts.filter(score__gte=(quiz.total_marks * Decimal('0.5'))).count()
            success_rate = (success_count / total_attempts * 100) if total_attempts else 0

            # 3️⃣ Score range distribution
            score_ranges = []
            for lower, upper in [(0, 20), (21, 40), (41, 60), (61, 80), (81, 100)]:
                count = attempts.filter(
                    score__gte=(quiz.total_marks * lower / 100),
                    score__lte=(quiz.total_marks * upper / 100)
                ).count()
                percent = (count / total_attempts * 100) if total_attempts else 0
                score_ranges.append((f"{lower}-{upper}%", count, round(percent, 1)))

            # 4️⃣ Top 10 students for this quiz
            top_students = (
                attempts.select_related('student')
                .order_by('-score')[:10]
            )

            # 5️⃣ Build quiz entry
            quiz_data.append({
                'quiz': quiz,
                'total_attempts': total_attempts,
                'avg_score': round(avg_score, 2),
                'success_rate': round(success_rate, 2),
                'score_ranges': score_ranges,
                'top_students': top_students,
            })

            # Update global stats
            total_quizzes += 1
            total_attempts_count += total_attempts
            total_score_sum += avg_score * total_attempts
            success_total += success_count

        # Add per-course summary
        course_data.append({
            'course': course,
            'quizzes': quiz_data,
        })

    # 6️⃣ Compute global analytics
    avg_score = (total_score_sum / total_attempts_count) if total_attempts_count else 0
    success_rate = (success_total / total_attempts_count * 100) if total_attempts_count else 0
    avg_attempts = (total_attempts_count / total_quizzes) if total_quizzes else 0

    # 7️⃣ Pass context to template
    context = {
        'total_courses': courses.count(),   # ✅ added for card display
        'course_data': course_data,
        'total_quizzes': total_quizzes,
        'avg_score': round(avg_score, 2),
        'success_rate': round(success_rate, 2),
        'avg_attempts': round(avg_attempts, 2),
    }
    return render(request, 'teacher_analytics.html', context)

from django.shortcuts import render, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.db.models import Avg, Max
from .models import Quiz, QuizAttempt, Course

@login_required
def quiz_analytics(request, quiz_id):
    # Get quiz without restricting to teacher, since course may be reassigned
    quiz = get_object_or_404(Quiz, id=quiz_id)

    # Ensure current teacher is assigned to the course related to this quiz
    if not Course.objects.filter(id=quiz.course.id, assigned_teacher=request.user).exists():
        return render(request, "access_denied.html", {"message": "You are not authorized to view this quiz analytics."})

    # All attempts for this quiz
    attempts = QuizAttempt.objects.filter(quiz=quiz)
    total_attempts = attempts.count()

    # --- Summary statistics ---
    if total_attempts > 0:
        avg_score = round(attempts.aggregate(avg=Avg('score'))['avg'] or 0, 1)
        max_score = round(attempts.aggregate(max=Max('score'))['max'] or 0, 1)
        passed_count = attempts.filter(score__gte=50).count()
        pass_rate = round((passed_count / total_attempts) * 100, 1)
    else:
        avg_score = max_score = pass_rate = 0

    # --- Top 5 students by score ---
    top_students = attempts.order_by('-score', 'submitted_at')[:5]


    context = {
        'quiz': quiz,
        'total_attempts': total_attempts,
        'avg_score': avg_score,
        'max_score': max_score,
        'pass_rate': pass_rate,
        'top_students': top_students,
    }

    return render(request, 'quiz_analytics.html', context)


from django.contrib.auth.decorators import login_required
from django.shortcuts import render
from django.db.models import F
from .models import QuizAttempt

@login_required
def student_results(request):
    user = request.user
    if not user.is_authenticated:
        return redirect('login')

    # Get all attempts for this student
    attempts = QuizAttempt.objects.select_related('quiz', 'quiz__course').filter(student=user)

    results = []

    for attempt in attempts:
        quiz = attempt.quiz
        total_marks = quiz.total_marks if hasattr(quiz, 'total_marks') else 100  # default 100
        percentage = round((attempt.score / total_marks) * 100, 2)

        # Rank calculation for this quiz
        all_attempts = QuizAttempt.objects.filter(quiz=quiz).order_by('-score', 'submitted_at')
        all_ids = list(all_attempts.values_list('id', flat=True))
        rank = all_ids.index(attempt.id) + 1  # +1 since index starts at 0
        total_students = len(all_ids)

        results.append({
            'quiz': quiz,
            'course': quiz.course.name,
            'score': attempt.score,
            'total_marks': total_marks,
            'percentage': percentage,
            'rank': rank,
            'total_students': total_students,
            'submitted_at': attempt.submitted_at,
        })

    return render(request, 'student_results.html', {'results': results})



from django.shortcuts import render, get_object_or_404, redirect
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from django.utils import timezone
from .models import QuizAttempt, QuizAnswer

@login_required
def quiz_result(request, attempt_id):
    # Get attempt safely
    attempt = get_object_or_404(QuizAttempt.objects.select_related("quiz", "student"), id=attempt_id)
    quiz = attempt.quiz

    # Permission check: Student can view only their own attempts
    if request.user.role == "Student" and attempt.student != request.user:
        messages.error(request, "You are not authorized to view this attempt.")
        return redirect('home')

    # (Optional) Teacher can only view attempts related to their assigned courses
    if request.user.role == "Teacher" and quiz.assigned_teacher != request.user:
        messages.error(request, "You are not authorized to view this quiz attempt.")
        return redirect('home')

    # Fetch all answers
    answers = QuizAnswer.objects.filter(attempt=attempt).select_related("question")

    correct = wrong = not_answered = 0
    question_answers = []

    for ans in answers:
        if ans.selected_option:
            if ans.is_correct:
                correct += 1
            else:
                wrong += 1
        else:
            not_answered += 1

        marks_obtained = ans.question.marks if ans.is_correct else 0
        question_answers.append({
            'question': ans.question,
            'student_answer': ans.selected_option,
            'is_correct': ans.is_correct,
            'marks_obtained': marks_obtained
        })

    total_questions = quiz.questions.count()
    total_score = float(attempt.score)
    max_score = total_questions * float(getattr(quiz, 'marks_per_question', 1))
    percentage = round((total_score / max_score) * 100, 2) if max_score > 0 else 0

    # Grade system
    if percentage >= 90:
        grade = 'A+'
    elif percentage >= 80:
        grade = 'A'
    elif percentage >= 70:
        grade = 'B'
    elif percentage >= 60:
        grade = 'C'
    elif percentage >= 50:
        grade = 'D'
    else:
        grade = 'F'

    context = {
        'quiz': quiz,
        'attempt': attempt,
        'question_answers': question_answers,
        'correct': correct,
        'wrong': wrong,
        'not_answered': not_answered,
        'percentage': percentage,
        'grade': grade,
    }

    return render(request, 'quiz_result.html', context)

# ------------------------------
# Attempt Details JSON (for modal)
# ------------------------------
from django.http import JsonResponse
from django.contrib.auth.decorators import login_required
from django.shortcuts import get_object_or_404
from .models import QuizAttempt, QuizAnswer

@login_required
def attempt_details_json(request, attempt_id):
    try:
        attempt = get_object_or_404(
            QuizAttempt.objects.select_related("quiz", "student"), id=attempt_id
        )

        # ✅ Permission checks
        if request.user.role == "Student" and attempt.student != request.user:
            return JsonResponse({"error": "Unauthorized access"}, status=403)
        if request.user.role == "Teacher" and attempt.quiz.teacher != request.user:
            return JsonResponse({"error": "Unauthorized access"}, status=403)

        answers = QuizAnswer.objects.filter(attempt=attempt).select_related("question")

        data = {
            "quiz_title": attempt.quiz.title,
            "score": float(attempt.score),
            "submitted_at": attempt.submitted_at.strftime("%d %b, %Y %H:%M")
            if attempt.submitted_at
            else "Not submitted",
            "questions": [],
        }

        for ans in answers:
            q = ans.question
            options = {
                "A": q.option_a,
                "B": q.option_b,
                "C": q.option_c,
                "D": q.option_d,
            }

            selected_text = options.get(ans.selected_option, "Not Answered")
            correct_text = options.get(q.correct_ans, "N/A")

            data["questions"].append({
                "question": q.question_text,
                "options": options,
                "selected_option": ans.selected_option,
                "selected_text": selected_text,
                "correct_option": q.correct_ans,
                "correct_text": correct_text,
                "is_correct": ans.is_correct,
            })

        return JsonResponse(data)

    except Exception as e:
        print("Error loading attempt details:", e)
        return JsonResponse({"error": "Failed to load details"}, status=500)
