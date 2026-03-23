import google.generativeai as genai
import json
import re
from django.conf import settings
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import io

# -------------------------
# Configure Gemini API Key
# -------------------------
genai.configure(api_key=settings.GEMINI_API_KEY)

# -------------------------
# 1️⃣ AI MCQ Generator
# -------------------------
def generate_mcqs_with_gemini(content, num_questions=5):
    """
    Generate MCQs using Gemini 2.5 models.
    Returns a list of dicts:
    [
      {
        "question_text": "...",
        "option_a": "...",
        "option_b": "...",
        "option_c": "...",
        "option_d": "...",
        "correct_ans": "A"
      }
    ]
    """
    if not content.strip():
        return []

    prompt = f"""
    Generate {num_questions} multiple-choice questions from the following study material.
    Each question must have 4 options (A, B, C, D) and specify the correct answer.
    Return strictly in JSON format like this:

    [
      {{
        "question_text": "...",
        "option_a": "...",
        "option_b": "...",
        "option_c": "...",
        "option_d": "...",
        "correct_ans": "A"
      }}
    ]

    Study Material:
    {content}
    """

    try:
        model = genai.GenerativeModel("models/gemini-2.5-flash")
        response = model.generate_content(prompt)
        text = response.text.strip()

        # Try parsing JSON
        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            # fallback using regex
            data = []
            questions = re.findall(r'"question_text"\s*:\s*"([^"]+)"', text)
            options_a = re.findall(r'"option_a"\s*:\s*"([^"]+)"', text)
            options_b = re.findall(r'"option_b"\s*:\s*"([^"]+)"', text)
            options_c = re.findall(r'"option_c"\s*:\s*"([^"]+)"', text)
            options_d = re.findall(r'"option_d"\s*:\s*"([^"]+)"', text)
            corrects = re.findall(r'"correct_ans"\s*:\s*"([A-D])"', text)

            for i in range(len(questions)):
                data.append({
                    "question_text": questions[i] if i < len(questions) else "",
                    "option_a": options_a[i] if i < len(options_a) else "",
                    "option_b": options_b[i] if i < len(options_b) else "",
                    "option_c": options_c[i] if i < len(options_c) else "",
                    "option_d": options_d[i] if i < len(options_d) else "",
                    "correct_ans": corrects[i] if i < len(corrects) else "A"
                })

        # Ensure exactly num_questions
        if len(data) > num_questions:
            data = data[:num_questions]
        elif len(data) < num_questions:
            for _ in range(num_questions - len(data)):
                data.append({
                    "question_text": "New Question",
                    "option_a": "",
                    "option_b": "",
                    "option_c": "",
                    "option_d": "",
                    "correct_ans": "A"
                })

        return data
    except Exception as e:
        print("❌ Gemini AI Error:", e)
        return []

# -------------------------
# 2️⃣ File Text Extraction
# -------------------------
def extract_text_from_file(uploaded_file):
    """
    Extract readable text from PDF, PPTX, or DOCX files.
    """
    text_content = ""
    filename = uploaded_file.name.lower()

    try:
        if filename.endswith(".pdf"):
            pdf_reader = PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text_content += page.extract_text() or ""

        elif filename.endswith((".ppt", ".pptx")):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"

        elif filename.endswith((".doc", ".docx")):
            doc = docx.Document(uploaded_file)
            for para in doc.paragraphs:
                text_content += para.text + "\n"

        else:
            # fallback for plain text
            text_content = uploaded_file.read().decode("utf-8", errors="ignore")

    except Exception as e:
        print("❌ File text extraction failed:", e)
        return ""

    return text_content.strip()
